from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
from langgraph.graph import StateGraph, END
from langgraph.graph.message import add_messages
from langgraph.checkpoint.memory import MemorySaver
from langgraph.prebuilt import ToolNode, tools_condition
from langchain_community.utilities import SerpAPIWrapper
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.tools import tool
from langchain_core.messages import HumanMessage, AIMessage
from dotenv import load_dotenv
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
from .teachers import anil_prompt, kavita_prompt, raghav_prompt, mary_prompt
from typing import TypedDict, List, Literal, Annotated, Optional, Dict, Any, AsyncGenerator
from pdfminer.high_level import extract_text
import os
import logging
from pptx import Presentation
from docx import Document
from cachetools import TTLCache
# Load environment variables FIRST
load_dotenv()

vector_stores = TTLCache(maxsize=1000, ttl=1800)

embeddings = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2",
    model_kwargs={'device': 'cpu'},
    encode_kwargs={'normalize_embeddings': True}
)

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Global storage for vector databases and conversation contexts
vector_stores: Dict[str, Chroma] = {}
conversation_contexts: Dict[str, Dict[str, Any]] = {}

class State(TypedDict):
    messages: Annotated[List, add_messages]
    teacher: Literal['Anil Deshmukh', 'Kavita Iyer', 'Raghav Sharma', 'Mary Fernandes']
    chain: Any
    pdf_path: Optional[str]
    user_id: str

def create_google_llm():
    """Create Google LLM with proper error handling"""
    try:
        api_key = os.getenv('GOOGLE_API_KEY')
        if not api_key:
            raise ValueError("GOOGLE_API_KEY environment variable is not set")
        
        llm = ChatGoogleGenerativeAI(
            model='gemini-2.0-flash',
            google_api_key=api_key,
            temperature=0.7
        )
        logger.info("Google LLM initialized successfully")
        return llm
    except Exception as e:
        logger.error(f"Failed to initialize Google LLM: {e}")
        raise

@tool
def search(query: str) -> str:
    """Search for information on the internet."""
    try:
        serpapi = SerpAPIWrapper()
        results = serpapi.run(query)
        return results
    except Exception as e:
        logger.error(f"Search error: {e}")
        return f"Search failed: {str(e)}"

# Initialize tools list
tools = [search]
tool_node = ToolNode(tools)

def extract_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_docx(docx_path):
    doc = Document(docx_path)
    full_text = []
    for para in doc.paragraphs:
        if para.text:
            full_text.append(para.text)
    return "\n".join(full_text)


async def prepare_pdf_rag(pdf_path: str, user_id: str) -> Chroma:
    # Check if this user already has a vectorstore
    existing = vector_stores.get(user_id)
    if existing and existing["file"] == pdf_path:
        # Reuse the same DB
        return existing["db"]
    
    # If new file or none exists, build fresh
    if pdf_path.endswith('.pdf'):
        content = extract_text(pdf_path)
    elif pdf_path.endswith('.txt'):
        with open(pdf_path, 'r') as f:
            content = f.read()
    elif pdf_path.endswith('.pptx'):
        content = extract_pptx(pdf_path)
    elif pdf_path.endswith('.docx'):
        content = extract_docx(pdf_path)
    
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200
    )
    chunks = splitter.split_text(content)
    
    vector_db = Chroma.from_texts(chunks, embeddings)
    
    # Store/replace
    vector_stores[user_id] = {"file": pdf_path, "db": vector_db}
    
    return vector_db

def initialise_teacher(state: State):
    """Initialize teacher with appropriate prompt and tools"""
    try:
        llm = create_google_llm()
        teacher = state['teacher']
        
        # Get teacher prompt
        if teacher == 'Anil Deshmukh':
            system_prompt = anil_prompt
        elif teacher == 'Kavita Iyer':
            system_prompt = kavita_prompt
        elif teacher == 'Raghav Sharma':
            system_prompt = raghav_prompt
        else:
            system_prompt = mary_prompt
        
        # Create a simpler prompt template
        prompt_template = ChatPromptTemplate.from_messages([
            ("system", system_prompt),
            ("user", "{input}")
        ])
        
        # Create the chain with the new template
        chain = prompt_template | llm.bind_tools(tools)
        
        return {"chain": chain}
        
    except Exception as e:
        logger.error(f"Error initializing teacher: {e}")
        raise

async def chat(state: State):
    """Modified chat handler with RAG support"""
    try:
        chain = state['chain']
        
        # Get the last message
        last_message = state['messages'][-1].content if state['messages'] else ""
        print(last_message) 
        if not last_message:
            return {"messages": [AIMessage(content="I didn't receive any message. Please try again.")]}

        # If PDF provided, use retriever to get relevant chunks
        if state.get('pdf_path'):
            vector_db = await prepare_pdf_rag(state['pdf_path'], state['user_id'])
            relevant_docs = vector_db.similarity_search(last_message, k=3)

            context = "\n\n".join([doc.page_content for doc in relevant_docs])
            
            # Format input with context
            input_text = f"Context from PDF:\n{context}\n\nQuestion: {last_message}"
        else:
            input_text = last_message

        # Invoke chain with simplified input
        try:
            response = await chain.ainvoke({"input": input_text})
            response = response.model_dump() if hasattr(response, 'model_dump') else None
            print(response)
            return {"messages": [AIMessage(content=response)]}
        except Exception as chain_error:
            logger.error(f"Chain invocation error: {chain_error}")
            return {"messages": [AIMessage(content=f"I encountered an error processing your request: {str(chain_error)}")]}

    except Exception as e:
        logger.error(f"Error in chat with RAG: {e}")
        return {"messages": [AIMessage(content=f"Sorry, I encountered an error: {str(e)}")]}

    
memory = MemorySaver()

# Build the graph
graph_builder = StateGraph(State)

# Add nodes
graph_builder.add_node("teacher", initialise_teacher)
graph_builder.add_node("chat", chat)
graph_builder.add_node("tools", tool_node)

graph_builder.add_conditional_edges("chat", tools_condition)
graph_builder.add_edge("teacher", "chat")
graph_builder.add_edge("tools", "chat")

# Set entry point
graph_builder.set_entry_point("teacher")

# Compile graph
try:
    agent = graph_builder.compile(checkpointer=memory)
    logger.info("Graph compiled successfully")
except Exception as e:
    logger.error(f"Error compiling graph: {e}")
    raise