from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
from langgraph.graph import StateGraph, END, START
from langgraph.graph.message import add_messages
from dotenv import load_dotenv
from typing import TypedDict, List, Literal, Annotated, Optional, Dict, Any, AsyncGenerator
from pdfminer.high_level import extract_text
from pydantic import BaseModel
from pptx import Presentation
from docx import Document
from sympy import content
import logging
import asyncio

load_dotenv()

# Setup logger for error reporting
logger = logging.getLogger("app_logger")
logger.setLevel(logging.DEBUG)  # or INFO in production

class Flashcards(BaseModel):
    question_1: str
    answer_1: str
    question_2: str
    answer_2: str
    question_3: str
    answer_3: str
    question_4: str
    answer_4: str
    question_5: str
    answer_5: str
    question_6: str
    answer_6: str
    question_7: str
    answer_7: str
    question_8: str
    answer_8: str
    question_9: str
    answer_9: str
    question_10: str
    answer_10: str

class Quiz(BaseModel):
    question_1: str
    options_1: List[str]
    answer_1: str
    question_2: str
    options_2: List[str]
    answer_2: str
    question_3: str
    options_3: List[str]
    answer_3: str
    question_4: str
    options_4: List[str]
    answer_4: str
    question_5: str
    options_5: List[str]
    answer_5: str
    question_6: str
    options_6: List[str]
    answer_6: str
    question_7: str
    options_7: List[str]
    answer_7: str
    question_8: str
    options_8: List[str]
    answer_8: str
    question_9: str
    options_9: List[str]
    answer_9: str
    question_10: str
    options_10: List[str]
    answer_10: str


class State(TypedDict):
    messages: Annotated[List, add_messages]
    teacher: Literal['Anil Deshmukh', 'Kavita Iyer', 'Raghav Sharma', 'Mary Fernandes']
    pdf_path: Optional[str]
    flashcards: Optional[Flashcards]
    quiz: Optional[Quiz]
    summarize: Optional[str]
    content: str
    result: any

def extract_pdf(pdf_path: str) -> str:
    try:
        text = extract_text(pdf_path)
        return text
    except Exception as e:
        logger.error(f"PDF extraction failed for {pdf_path}: {str(e)}")
        return ""

def extract_pptx(pptx_path: str) -> str:
    try:
        prs = Presentation(pptx_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        logger.error(f"PPTX extraction failed for {pptx_path}: {str(e)}")
        return ""

def extract_docx(docx_path: str) -> str:
    try:
        doc = Document(docx_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text:
                full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"DOCX extraction failed for {docx_path}: {str(e)}")
        return ""

def extract_file(state: State) -> str:
    try:
        content = ""
        if state.get('pdf_path'):
            pdf_path = state.get('pdf_path')
            if pdf_path.endswith('.pdf'):
                content = extract_pdf(pdf_path)
            elif pdf_path.endswith('.txt'):
                try:
                    with open(pdf_path, 'r') as file:
                        content = file.read()
                except Exception as e:
                    logger.error(f"Reading TXT file failed: {pdf_path} : {str(e)}")
                    content = ""
            elif pdf_path.endswith('.pptx'):
                content = extract_pptx(pdf_path)
            elif pdf_path.endswith('.docx'):
                content = extract_docx(pdf_path)
            else:
                logger.warning(f"Unsupported file extension for extraction: {pdf_path}")
                content = ""
        else:
            # If no file, fallback to last message in messages list safely
            messages = state.get('messages', [])
            if messages and isinstance(messages, list):
                content = messages[-1] if messages else ""
            else:
                logger.warning("State messages missing or invalid when extracting content.")
                content = ""

        return {"content": content}
    except Exception as e:
        logger.error(f"Unhandled exception in extract_file: {str(e)}")
        return ""

async def summarize(state: State) -> Dict[str, Optional[str]]:
    try:
        llm = ChatGoogleGenerativeAI(
            model='gemini-2.0-flash',
            temperature=0.7
        )

        content = state['content']
        if not content:
            raise ValueError("No content available to summarize.")

        prompt = ChatPromptTemplate.from_messages([
            ("system", """
                You are a summarizer. Your task is to take the provided text or query 
                (which may come from a user message or an uploaded file) and generate a 
                clear, concise summary. 
                Focus on the main ideas, key details, and important context. 
                Keep the language simple and easy to understand.
            """),
            ("human", "{content}")
        ])

        chain = prompt | llm
        result = await chain.ainvoke({"content": content})
        return {"summarize": result.content}

    except Exception as e:
        logger.error(f"Error in summarize: {str(e)}")
        return {"summarize": None}

async def generate_quiz(state: State) -> Dict[str, Optional[Dict]]:
    try:
        llm = ChatGoogleGenerativeAI(
            model='gemini-2.0-flash',
            temperature=0.7
        ).with_structured_output(Quiz)

        content = state['content']
        if not content:
            raise ValueError("No content available to generate quiz.")

        prompt = ChatPromptTemplate.from_messages([
            ("system", """
                You are an expert educator. The user will provide a topic, and you must generate questions on that topic using Bloom’s Taxonomy. 
                Create questions at different cognitive levels: Remember, Understand, Apply, Analyze, Evaluate, and Create. 
                Label each question with its Bloom’s level. Provide 2–3 questions per level. 
                Ensure progression from simple factual recall to higher-order critical thinking and creativity.
            """),
            ("human", "{content}")
        ])

        chain = prompt | llm
        result = await chain.ainvoke({"content": content})

        flashcard_dict = result.model_dump() if hasattr(result, 'model_dump') else None
        return {"quiz": flashcard_dict}

    except Exception as e:
        logger.error(f"Error in generate_quiz: {str(e)}")
        return {"quiz": None}

async def generate_flashcards(state: State) -> Dict[str, Optional[Dict]]:
    try:
        llm = ChatGoogleGenerativeAI(
            model='gemini-2.0-flash',
            temperature=0.7
        ).with_structured_output(Flashcards)

        content = state['content']
        if not content:
            raise ValueError("No content available to generate flashcards.")

        prompt = ChatPromptTemplate.from_messages([
            ("system", """You are a flash card generator. Your job is to create 10 question-answer pairs from the provided content.
                Focus on key concepts and important details. Make questions clear and concise."""),
            ("human", "Generate 10 flashcards from this content: {content}")
        ])

        chain = prompt | llm
        result = await chain.ainvoke({"content": content})

        flashcard_dict = result.model_dump() if hasattr(result, 'model_dump') else None
        return {"flashcards": flashcard_dict}

    except Exception as e:
        logger.error(f"Error in generate_flashcards: {str(e)}")
        return {"flashcards": None}

async def chat(state: State) -> Dict[str, Optional[Any]]:
    try:
        return {
            "result": {
                "flashcards": state.get("flashcards"),
                "quiz": state.get("quiz"),
                "summarize": state.get("summarize")
            }
        }
    except Exception as e:
        logger.error(f"Error in chat aggregation: {str(e)}")
        return {"result": None}

# Build graph with error handling integrated
graph_builder = StateGraph(State)
graph_builder.add_node("quiz", generate_quiz)
graph_builder.add_node("flashcards", generate_flashcards)
graph_builder.add_node("summarize", summarize)
graph_builder.add_node("extract", extract_file)
graph_builder.set_entry_point("extract")
graph_builder.add_node("chat", chat)
graph_builder.add_edge("extract", "quiz")
graph_builder.add_edge("extract", "flashcards")
graph_builder.add_edge("extract", "summarize")
graph_builder.add_edge("quiz", "chat")
graph_builder.add_edge("flashcards", "chat")
graph_builder.add_edge("summarize", "chat")
graph_builder.set_finish_point("chat")

graph = graph_builder.compile()

print(graph.get_graph().draw_mermaid())
