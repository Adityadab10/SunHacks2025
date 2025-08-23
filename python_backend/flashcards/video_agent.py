from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
from langgraph.graph import StateGraph, END, START
from langgraph.graph.message import add_messages
from dotenv import load_dotenv
from typing import TypedDict, List, Literal, Annotated, Optional, Dict, Any, AsyncGenerator
from pdfminer.high_level import extract_text
from pydantic import BaseModel
from pptx import Presentation
from docx import Document
from sympy import content
import logging

load_dotenv()

# Setup logger for error reporting
logger = logging.getLogger("app_logger")
logger.setLevel(logging.DEBUG)  # or INFO in production

class ImportantPoint(BaseModel):
    points: List[str]

class State(TypedDict):
    messages: Annotated[List, add_messages]
    important_points: ImportantPoint
    pdf_path: Optional[str]
    content: str

def extract_pdf(pdf_path: str) -> str:
    try:
        text = extract_text(pdf_path)
        return text
    except Exception as e:
        logger.error(f"PDF extraction failed for {pdf_path}: {str(e)}")
        return ""

def extract_pptx(pptx_path: str) -> str:
    try:
        prs = Presentation(pptx_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        logger.error(f"PPTX extraction failed for {pptx_path}: {str(e)}")
        return ""

def extract_docx(docx_path: str) -> str:
    try:
        doc = Document(docx_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text:
                full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"DOCX extraction failed for {docx_path}: {str(e)}")
        return ""

def extract_file(state: State) -> str:
    try:
        content = ""
        if state.get('pdf_path'):
            pdf_path = state.get('pdf_path')
            if pdf_path.endswith('.pdf'):
                content = extract_pdf(pdf_path)
            elif pdf_path.endswith('.txt'):
                try:
                    with open(pdf_path, 'r') as file:
                        content = file.read()
                except Exception as e:
                    logger.error(f"Reading TXT file failed: {pdf_path} : {str(e)}")
                    content = ""
            elif pdf_path.endswith('.pptx'):
                content = extract_pptx(pdf_path)
            elif pdf_path.endswith('.docx'):
                content = extract_docx(pdf_path)
            else:
                logger.warning(f"Unsupported file extension for extraction: {pdf_path}")
                content = ""
        else:
            # If no file, fallback to last message in messages list safely
            messages = state.get('messages', [])
            if messages and isinstance(messages, list):
                content = messages[-1] if messages else ""
            else:
                logger.warning("State messages missing or invalid when extracting content.")
                content = ""

        return {"content": content}
    except Exception as e:
        logger.error(f"Unhandled exception in extract_file: {str(e)}")
        return ""

async def generate_important(state: State) -> Dict[str, Optional[Dict]]:
    try:
        llm = ChatGoogleGenerativeAI(
            model='gemini-2.0-flash',
            temperature=0.7
        ).with_structured_output(ImportantPoint)

        content = state['content']
        if not content:
            raise ValueError("No content available to generate important points.")

        prompt = ChatPromptTemplate.from_messages([
            ("system", """You are an important points generator. Your job is to create key points from the provided content.
                Focus on key concepts and important details. Make points clear and concise."""),
            ("human", "Generate important points from this content: {content}")
        ])

        chain = prompt | llm
        result = await chain.ainvoke({"content": content})
        print(f"result generated: {result}")
        return {"result": result.points}

    except Exception as e:
        logger.error(f"Error in generate_important: {str(e)}")
        return {"important_points": None}

# Build graph with error handling integrated
graph_builder = StateGraph(State)
graph_builder.add_node("imp", generate_important)
graph_builder.add_node("extract", extract_file)
graph_builder.set_entry_point("extract")
graph_builder.add_edge("extract", "imp")
graph_builder.set_finish_point("imp")

graph_imp = graph_builder.compile()

print(graph_imp.get_graph().draw_mermaid())
